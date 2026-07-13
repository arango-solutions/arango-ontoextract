"""Document parsing, semantic chunking, and duplicate detection.

Supports:

* **PDF**  via ``pymupdf`` (``fitz``)
* **DOCX** via ``python-docx``
* **PPTX** via ``python-pptx``
* **DOC**  (legacy Word binary, pre-2007) via a LibreOffice subprocess
  that converts ``.doc`` -> ``.docx`` and then reuses :func:`parse_docx`.
  Requires LibreOffice/soffice on PATH (``brew install --cask libreoffice``
  on macOS, ``apt install libreoffice-core`` on Debian). Fails loudly
  with a clear error if not present.
* **Markdown** plain-text parser
"""

from __future__ import annotations

import hashlib
import io
import logging
import os
import re
import shutil
import subprocess
import tempfile
from collections.abc import Iterable
from dataclasses import dataclass, field
from typing import Any

import fitz
import tiktoken
from docx import Document
from pptx import Presentation

from app.config import settings
from app.services.visual_extraction import (
    VisualAsset,
    VisualExtractionDiagnostics,
    collect_pptx_visual_assets,
    format_section_chunk_text,
    get_caption_provider,
    visual_placeholder_line,
)

log = logging.getLogger(__name__)

# Subprocess timeout for the LibreOffice .doc -> .docx conversion.
# 60s handles large legacy decks; the convert pass itself rarely exceeds 5s.
_LIBREOFFICE_TIMEOUT_SECONDS = 60

_TIKTOKEN_MODEL = "cl100k_base"


@dataclass
class Section:
    heading: str
    text: str
    page_number: int | None = None
    #: Indexes into ``ParsedDocument.visual_diagnostics.assets`` for
    #: assets attached to this section. Used by chunking (IMG.5) to
    #: propagate per-chunk visual context.
    visual_asset_indexes: list[int] = field(default_factory=list)
    #: Speaker notes for this section (CH.2). For decks this is the slide's
    #: notes-pane text (without the ``[Notes]`` marker). The slide-aware
    #: chunker (``_chunk_deck``) emits it as a distinct chunk linked to the
    #: slide; the legacy chunker folds it into the section body with a
    #: ``[Notes]`` marker (preserving pre-Stream-17 behavior).
    notes: str = ""


@dataclass
class ParsedDocument:
    sections: list[Section] = field(default_factory=list)
    title: str = ""
    author: str = ""
    page_count: int = 0
    format: str = ""
    visual_diagnostics: VisualExtractionDiagnostics = field(
        default_factory=VisualExtractionDiagnostics
    )


@dataclass
class Chunk:
    text: str
    chunk_index: int
    source_page: int | None
    section_heading: str
    token_count: int
    #: One of ``"text"`` (no visual placeholders), ``"visual"`` (only
    #: visual placeholders / scanned-page markers), or ``"mixed"``.
    #: Used by the visual-aware extraction strategy (IMG.6) and by
    #: curators inspecting why a chunk's content looks unusual.
    chunk_kind: str = "text"
    #: Per-chunk projection of the visual assets that contributed to
    #: this chunk's text. Mirrors ``VisualAsset.to_dict()`` shape.
    visual_assets: list[dict[str, Any]] = field(default_factory=list)
    #: Source document format (``"pptx"`` | ``"pdf"`` | ``"docx"`` |
    #: ``"markdown"`` | ``"unknown"``), copied from
    #: ``ParsedDocument.format``. Persisted on stored chunks (CH.1) so the
    #: extraction strategy selector can detect decks in production instead
    #: of relying on visual markers alone (FR-1.17).
    doc_format: str = ""
    #: Role of the chunk within its source section (CH.2). ``"body"`` for
    #: ordinary slide/section content; ``"notes"`` for a speaker-notes chunk
    #: emitted separately by the slide-aware deck chunker. Persisted only
    #: when non-default so legacy/non-deck storage is byte-identical.
    chunk_role: str = "body"
    #: When a single slide exceeds ``chunk_max_tokens`` the deck chunker
    #: splits it across ``slide_parts`` chunks; ``slide_part`` is this
    #: chunk's 0-based index within that slide (CH.2 "record the split").
    #: ``slide_parts == 1`` means the slide was not split. Defaults keep
    #: non-deck chunks byte-identical.
    slide_part: int = 0
    slide_parts: int = 1
    #: Topic-unit id (CH.3): consecutive slides that continue one topic
    #: (repeated / ``(cont'd)`` titles) share a ``topic_unit`` so the
    #: extractor batches by topic unit instead of raw chunk index. ``None``
    #: when topic grouping does not apply (e.g. non-deck documents).
    topic_unit: int | None = None


def compute_file_hash(content: bytes) -> str:
    """SHA-256 hex digest of raw file bytes."""
    return hashlib.sha256(content).hexdigest()


def _token_count(text: str) -> int:
    enc = tiktoken.get_encoding(_TIKTOKEN_MODEL)
    return len(enc.encode(text))


# ---------------------------------------------------------------------------
# Parsers
# ---------------------------------------------------------------------------


def parse_pdf(file_bytes: bytes) -> ParsedDocument:
    """Extract sections from a PDF using pymupdf (fitz)."""
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    parsed = ParsedDocument(page_count=len(doc), format="pdf")
    emit_placeholders = (
        settings.visual_extraction_enabled and settings.visual_extraction_placeholders
    )

    metadata = doc.metadata or {}
    parsed.title = metadata.get("title", "") or ""
    parsed.author = metadata.get("author", "") or ""

    # ``TEXT_PRESERVE_IMAGES`` is required for image-only pages: without
    # it, ``get_text("dict")`` returns zero blocks for a page whose only
    # content is an XObject image, and the scanned-page branch below
    # never fires (IMG.8 regression: scanned PDFs silently produced no
    # placeholders and ``scanned_page_count == 0``).
    pdf_text_flags = fitz.TEXT_PRESERVE_WHITESPACE | fitz.TEXT_PRESERVE_IMAGES
    for page_num, page in enumerate(doc, start=1):
        blocks = page.get_text("dict", flags=pdf_text_flags)["blocks"]
        current_heading = ""
        current_text_parts: list[str] = []
        image_blocks_on_page = 0
        page_section_start = len(parsed.sections)

        for block in blocks:
            if block.get("type") != 0:
                image_blocks_on_page += 1
                continue
            for line in block.get("lines", []):
                line_text = "".join(span["text"] for span in line.get("spans", []))
                if not line_text.strip():
                    continue

                max_size = max(
                    (span.get("size", 0) for span in line.get("spans", [])),
                    default=0,
                )
                is_bold = any(
                    "bold" in (span.get("font", "").lower()) for span in line.get("spans", [])
                )

                if max_size >= 14 or (is_bold and max_size >= 12):
                    if current_text_parts:
                        parsed.sections.append(
                            Section(
                                heading=current_heading,
                                text="\n".join(current_text_parts).strip(),
                                page_number=page_num,
                            )
                        )
                        current_text_parts = []
                    current_heading = line_text.strip()
                else:
                    current_text_parts.append(line_text)

        if current_text_parts:
            parsed.sections.append(
                Section(
                    heading=current_heading,
                    text="\n".join(current_text_parts).strip(),
                    page_number=page_num,
                )
            )

        if settings.visual_extraction_enabled and image_blocks_on_page:
            visual_lines: list[str] = []
            new_asset_indexes: list[int] = []
            for idx in range(1, image_blocks_on_page + 1):
                parsed.visual_diagnostics.register_asset(
                    VisualAsset(
                        page_number=page_num,
                        asset_index=idx,
                        asset_type="image_block",
                        method="placeholder",
                    )
                )
                new_asset_indexes.append(parsed.visual_diagnostics.visual_asset_count - 1)
                if emit_placeholders:
                    visual_lines.append(
                        visual_placeholder_line(
                            page_number=page_num,
                            asset_index=idx,
                            asset_type="image_block",
                            doc_format="pdf",
                        )
                    )

            page_has_text = len(parsed.sections) > page_section_start
            if not page_has_text and image_blocks_on_page:
                parsed.visual_diagnostics.scanned_page_count += 1
                if emit_placeholders:
                    visual_lines.insert(
                        0,
                        f"[Scanned or image-only page {page_num}: OCR not configured]",
                    )
                parsed.sections.append(
                    Section(
                        heading="",
                        text="\n".join(visual_lines),
                        page_number=page_num,
                        visual_asset_indexes=new_asset_indexes,
                    )
                )
            elif visual_lines:
                target = parsed.sections[-1]
                extra = "\n".join(visual_lines)
                target.text = f"{target.text}\n\n{extra}".strip() if target.text else extra
                target.visual_asset_indexes.extend(new_asset_indexes)

    doc.close()

    if not parsed.sections:
        full_text = ""
        reopened = fitz.open(stream=file_bytes, filetype="pdf")
        for page in reopened:
            full_text += page.get_text() + "\n"
        reopened.close()
        if full_text.strip():
            parsed.sections.append(Section(heading="", text=full_text.strip(), page_number=1))

    return parsed


def parse_docx(file_bytes: bytes) -> ParsedDocument:
    """Extract sections from a DOCX using python-docx."""
    doc = Document(io.BytesIO(file_bytes))

    parsed = ParsedDocument()
    core = doc.core_properties
    parsed.title = core.title or ""
    parsed.author = core.author or ""
    parsed.format = "docx"

    current_heading = ""
    current_text_parts: list[str] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        if para.style and para.style.name and para.style.name.startswith("Heading"):
            if current_text_parts:
                parsed.sections.append(
                    Section(heading=current_heading, text="\n".join(current_text_parts))
                )
                current_text_parts = []
            current_heading = text
        else:
            current_text_parts.append(text)

    if current_text_parts:
        parsed.sections.append(Section(heading=current_heading, text="\n".join(current_text_parts)))

    return parsed


def parse_pptx(file_bytes: bytes) -> ParsedDocument:
    """Extract sections from a PPTX (PowerPoint) using python-pptx.

    Each slide becomes one :class:`Section` so downstream chunking and
    provenance line up with what users intuitively expect ("page 7 of
    that deck"):

    * ``page_number`` = 1-based slide index (matches PDF semantics).
    * ``heading`` = slide title placeholder if present, else "".
    * ``text`` = concatenation of every other text-bearing shape on the
      slide (body placeholders, text boxes, table cells, grouped shapes
      flattened recursively).
    * ``notes`` = the slide's speaker-notes text (CH.2). It is carried on
      the :class:`Section` rather than folded into ``text`` so the
      slide-aware chunker can emit it as a distinct chunk linked to the
      slide; the legacy chunker re-folds it into the body with a
      ``[Notes]`` marker.

    Empty slides (no extractable text and no notes) are dropped to keep
    chunk counts honest.
    """
    prs = Presentation(io.BytesIO(file_bytes))
    parsed = ParsedDocument(page_count=len(prs.slides), format="pptx")
    emit_placeholders = (
        settings.visual_extraction_enabled and settings.visual_extraction_placeholders
    )
    caption_provider = (
        get_caption_provider(settings.visual_caption_provider)
        if settings.visual_extraction_enabled
        else None
    )

    core = prs.core_properties
    parsed.title = (core.title or "") if core else ""
    parsed.author = (core.author or "") if core else ""

    for slide_index, slide in enumerate(prs.slides, start=1):
        heading = _pptx_slide_title(slide)
        body_parts = _pptx_collect_text(slide.shapes, exclude_title=True)
        asset_count_before = parsed.visual_diagnostics.visual_asset_count

        if settings.visual_extraction_enabled:
            body_parts.extend(
                collect_pptx_visual_assets(
                    slide.shapes,
                    slide_index=slide_index,
                    diagnostics=parsed.visual_diagnostics,
                    emit_placeholders=emit_placeholders,
                    caption_provider=caption_provider,
                    max_caption_calls=settings.visual_caption_max_assets_per_doc,
                )
            )
        slide_asset_indexes = list(
            range(asset_count_before, parsed.visual_diagnostics.visual_asset_count)
        )

        # Speaker notes: powerful provenance signal, often where the
        # actual narrative lives in survey decks.
        notes_text = ""
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text or ""
                notes_text = notes.strip()
        except Exception:
            # Don't let a malformed notes pane kill the whole parse.
            log.warning("pptx parse: failed to read notes for slide %d", slide_index, exc_info=True)

        body = "\n".join(p for p in body_parts if p).strip()
        if not body and not heading and not notes_text:
            continue

        parsed.sections.append(
            Section(
                heading=heading,
                text=body,
                page_number=slide_index,
                visual_asset_indexes=slide_asset_indexes,
                notes=notes_text,
            )
        )

    return parsed


def _pptx_slide_title(slide: object) -> str:
    """Return the slide's title-placeholder text if any, else ""."""
    try:
        title_shape = slide.shapes.title  # type: ignore[attr-defined]
    except Exception:
        return ""
    if title_shape is None:
        return ""
    text = getattr(title_shape, "text", "") or ""
    return text.strip()


def _pptx_collect_text(
    shapes: Iterable[Any],
    *,
    exclude_title: bool,
) -> list[str]:
    """Walk a python-pptx shape collection (incl. groups + tables) for text.

    ``exclude_title`` skips the slide title shape so we don't double-count
    it (the caller already extracted it as the section heading).
    """
    out: list[str] = []
    for shape in shapes:
        if exclude_title and getattr(shape, "is_placeholder", False):
            ph = getattr(shape, "placeholder_format", None)
            # Title placeholder idx is 0 in OpenXML.
            if ph is not None and getattr(ph, "idx", None) == 0:
                continue

        # Grouped shapes -> recurse.
        if getattr(shape, "shape_type", None) == 6:  # MSO_SHAPE_TYPE.GROUP
            out.extend(_pptx_collect_text(shape.shapes, exclude_title=False))
            continue

        # Tables: row by row, cell by cell.
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                row_text = " | ".join((cell.text or "").strip() for cell in row.cells).strip(" |")
                if row_text:
                    out.append(row_text)
            continue

        # Plain text frames.
        if getattr(shape, "has_text_frame", False):
            text = (shape.text_frame.text or "").strip()
            if text:
                out.append(text)
    return out


def parse_doc(file_bytes: bytes) -> ParsedDocument:
    """Extract sections from a legacy ``.doc`` (Word 97-2003 binary) file.

    Strategy: shell out to LibreOffice in headless mode to convert the
    binary ``.doc`` to ``.docx``, then reuse :func:`parse_docx` so we
    keep one code path for Word styles + heading detection.

    LibreOffice / ``soffice`` must be on PATH. We probe at call-time
    rather than import-time so the rest of the module still imports
    on a host that has not yet installed it.

    Raises
    ------
    RuntimeError
        When LibreOffice is not installed or the conversion fails. The
        error message tells the operator exactly what to install.
    """
    soffice = _find_libreoffice()
    if soffice is None:
        raise RuntimeError(
            "Cannot parse legacy .doc files: LibreOffice (soffice) is not "
            "installed. Install it with `brew install --cask libreoffice` "
            "(macOS) or `apt install libreoffice-core` (Debian/Ubuntu) and "
            "retry. Alternatively, convert the file to .docx and re-upload."
        )

    with tempfile.TemporaryDirectory(prefix="aoe_doc_") as tmpdir:
        in_path = os.path.join(tmpdir, "in.doc")
        out_dir = os.path.join(tmpdir, "out")
        os.makedirs(out_dir, exist_ok=True)
        with open(in_path, "wb") as fh:
            fh.write(file_bytes)

        try:
            result = subprocess.run(
                [
                    soffice,
                    "--headless",
                    "--convert-to",
                    "docx",
                    "--outdir",
                    out_dir,
                    in_path,
                ],
                capture_output=True,
                timeout=_LIBREOFFICE_TIMEOUT_SECONDS,
                check=False,
            )
        except subprocess.TimeoutExpired as exc:
            raise RuntimeError(
                f"LibreOffice .doc -> .docx conversion timed out after "
                f"{_LIBREOFFICE_TIMEOUT_SECONDS}s. The file may be malformed "
                f"or unusually large."
            ) from exc

        if result.returncode != 0:
            stderr = (result.stderr or b"").decode("utf-8", errors="replace").strip()
            raise RuntimeError(
                f"LibreOffice conversion failed (exit={result.returncode}): "
                f"{stderr or '<no stderr>'}"
            )

        # The converted file lands in out_dir as <basename>.docx.
        candidates = [f for f in os.listdir(out_dir) if f.endswith(".docx")]
        if not candidates:
            raise RuntimeError(
                "LibreOffice conversion produced no .docx output. Stderr: "
                + (result.stderr or b"").decode("utf-8", errors="replace").strip()
            )

        with open(os.path.join(out_dir, candidates[0]), "rb") as fh:
            converted_bytes = fh.read()

    return parse_docx(converted_bytes)


def _find_libreoffice() -> str | None:
    """Locate the LibreOffice headless binary, or return None if absent.

    Checks PATH for ``soffice`` and ``libreoffice``, then the standard
    macOS install location. Order matters: PATH first so an admin can
    pin a specific version.
    """
    for name in ("soffice", "libreoffice"):
        path = shutil.which(name)
        if path:
            return path
    mac_default = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    if os.path.isfile(mac_default) and os.access(mac_default, os.X_OK):
        return mac_default
    return None


_MD_HEADING_RE = re.compile(r"^(#{1,6})\s+(.*)", re.MULTILINE)


def parse_markdown(text: str) -> ParsedDocument:
    """Extract sections from Markdown text using heading boundaries."""
    parsed = ParsedDocument(format="markdown")

    lines = text.split("\n")
    first_heading = ""
    for line in lines:
        m = _MD_HEADING_RE.match(line)
        if m:
            first_heading = m.group(2).strip()
            break
    parsed.title = first_heading

    current_heading = ""
    current_text_parts: list[str] = []

    for line in lines:
        m = _MD_HEADING_RE.match(line)
        if m:
            if current_text_parts:
                parsed.sections.append(
                    Section(heading=current_heading, text="\n".join(current_text_parts).strip())
                )
                current_text_parts = []
            current_heading = m.group(2).strip()
        else:
            current_text_parts.append(line)

    if current_text_parts:
        body = "\n".join(current_text_parts).strip()
        if body:
            parsed.sections.append(Section(heading=current_heading, text=body))

    return parsed


# ---------------------------------------------------------------------------
# Chunking
# ---------------------------------------------------------------------------


def _split_into_paragraphs(text: str) -> list[str]:
    """Split text at blank-line boundaries, keeping non-empty paragraphs."""
    paragraphs = re.split(r"\n\s*\n", text)
    return [p.strip() for p in paragraphs if p.strip()]


# CH.3: continuation markers on slide titles ("Financials (cont'd)",
# "Roadmap - continued", "Agenda 2") signal that a slide continues the
# previous slide's topic. We strip them to a base title and group
# consecutive slides sharing a base title into one "topic unit".
_CONT_SUFFIX_RE = re.compile(
    r"\s*[\(\[]?\s*(?:cont(?:inued|'d|d)?\.?)\s*[\)\]]?\s*$",
    re.IGNORECASE,
)
_TRAILING_NUM_RE = re.compile(r"\s*[\(\[]?\s*\d+\s*[\)\]]?\s*$")


def categorize_document(parsed: ParsedDocument) -> str:
    """Classify a parsed document into a chunk-strategy category (CH.4).

    Runs **before** chunking (FR-1.16) so the category selects the chunk
    strategy (boundary rules) rather than being inferred post-hoc from
    chunk text -- which is what :mod:`app.extraction.agents.strategy` still
    does for prompt/batch/pass selection at extraction time.

    Today the actionable distinction is *deck* vs. everything else: a deck
    (currently PPTX) routes to the slide-boundary-preserving chunker
    (:func:`_chunk_deck`). ``"narrative"`` is returned for all other
    formats and routes to the paragraph chunker (:func:`_chunk_prose`),
    which is byte-identical to the pre-Stream-17 behavior. The
    ``tabular`` / ``technical`` categories are intentionally *not* split
    out here yet -- those still rely on chunk-text heuristics in the
    strategy selector; promoting them pre-chunk is a follow-up that this
    hook makes cheap.
    """
    fmt = (parsed.format or "").lower()
    if fmt == "pptx":
        return "deck"
    return "narrative"


def chunk_document(
    parsed: ParsedDocument,
    max_tokens: int | None = None,
    *,
    category: str | None = None,
    overlap_tokens: int | None = None,
) -> list[Chunk]:
    """Chunk a parsed document at section / paragraph (or slide) boundaries.

    Each chunk respects ``max_tokens`` (counted via tiktoken
    ``cl100k_base``); when ``max_tokens`` is ``None`` the configured
    ``settings.chunk_max_tokens`` is used (CH.4 -- the size is no longer a
    hardcoded constant). Chunks preserve source page and section heading
    metadata.

    The ``category`` (see :func:`categorize_document`) selects the chunk
    strategy (CH.4 "categorize-then-chunk ordering"). Decks route to the
    slide-boundary-preserving chunker (CH.2/CH.3): a slide is never merged
    with another slide, is split only when it exceeds ``max_tokens`` (the
    split is recorded on ``slide_part`` / ``slide_parts``), speaker notes
    become a distinct ``chunk_role="notes"`` chunk linked to their slide,
    and continuation slides are grouped into ``topic_unit`` s. Deck-aware
    chunking can be disabled via ``settings.chunk_slide_aware`` (kill
    switch), in which case decks fall back to the paragraph chunker.

    ``overlap_tokens`` (default ``settings.chunk_overlap_tokens``) repeats
    trailing paragraphs -- up to that token budget -- at the head of the
    next chunk from the *same* section/slide; ``0`` (the default) yields no
    overlap and is byte-identical to the pre-Stream-17 chunker.

    IMG.5: chunks carry ``chunk_kind`` (``"text"`` | ``"visual"`` |
    ``"mixed"``) and ``visual_assets`` (the projection of visual assets
    that fell in the source section) so downstream prompts can label
    visual context distinctly without re-scanning the chunk text.
    """
    resolved_max = settings.chunk_max_tokens if max_tokens is None else max_tokens
    resolved_overlap = settings.chunk_overlap_tokens if overlap_tokens is None else overlap_tokens
    # Clamp overlap so it can never consume the whole window (which would
    # duplicate content without making forward progress).
    resolved_overlap = max(0, min(resolved_overlap, resolved_max // 2))
    resolved_category = category or categorize_document(parsed)
    doc_format = parsed.format or "unknown"

    if resolved_category == "deck" and settings.chunk_slide_aware:
        chunks, origin = _chunk_deck(parsed, resolved_max, resolved_overlap, doc_format)
    else:
        chunks, origin = _chunk_prose(parsed, resolved_max, resolved_overlap, doc_format)

    _backfill_visual_context(parsed, chunks, origin)
    return chunks


def _split_long_paragraph(para: str, max_tokens: int) -> list[str]:
    """Split a single oversized paragraph into <= ``max_tokens`` word runs."""
    pieces: list[str] = []
    sub_parts: list[str] = []
    sub_tokens = 0
    for word in para.split():
        word_tokens = _token_count(word + " ")
        if sub_tokens + word_tokens > max_tokens and sub_parts:
            pieces.append(" ".join(sub_parts))
            sub_parts = []
            sub_tokens = 0
        sub_parts.append(word)
        sub_tokens += word_tokens
    if sub_parts:
        pieces.append(" ".join(sub_parts))
    return pieces


def _overlap_tail(parts: list[str], budget: int) -> tuple[list[str], int]:
    """Trailing paragraphs (newest-last) whose token sum stays <= ``budget``."""
    if budget <= 0:
        return [], 0
    carry: list[str] = []
    total = 0
    for part in reversed(parts):
        part_tokens = _token_count(part)
        if total + part_tokens > budget:
            break
        carry.insert(0, part)
        total += part_tokens
    return carry, total


def _pack_paragraphs(
    paragraphs: list[str],
    max_tokens: int,
    overlap_tokens: int,
) -> list[str]:
    """Pack paragraphs into <= ``max_tokens`` text pieces.

    With ``overlap_tokens == 0`` this reproduces the pre-Stream-17 packing
    exactly (paragraph-boundary packing; oversized paragraphs word-split;
    no overlap). With a positive overlap, trailing paragraphs up to the
    overlap budget are repeated at the head of the following piece.
    """
    pieces: list[str] = []
    current_parts: list[str] = []
    current_tokens = 0

    def flush(*, carry_overlap: bool) -> None:
        nonlocal current_parts, current_tokens
        if not current_parts:
            return
        pieces.append("\n\n".join(current_parts))
        if carry_overlap and overlap_tokens > 0:
            current_parts, current_tokens = _overlap_tail(current_parts, overlap_tokens)
        else:
            current_parts, current_tokens = [], 0

    for para in paragraphs:
        para_tokens = _token_count(para)

        if para_tokens > max_tokens:
            # Emit whatever is buffered (no overlap into a word-split run),
            # then hard-split the oversized paragraph word-by-word.
            flush(carry_overlap=False)
            pieces.extend(_split_long_paragraph(para, max_tokens))
            continue

        if current_tokens + para_tokens > max_tokens and current_parts:
            flush(carry_overlap=True)

        current_parts.append(para)
        current_tokens += para_tokens

    flush(carry_overlap=False)
    return pieces


def _section_paragraphs(section: Section, doc_format: str, *, fold_notes: bool) -> list[str]:
    """Formatted, paragraph-split text for a section's body.

    ``fold_notes`` re-appends the section's speaker notes (with a
    ``[Notes]`` marker) into the body before formatting -- the legacy
    behavior used by the paragraph chunker so decks are unchanged when
    slide-aware chunking is disabled.
    """
    body = section.text
    if fold_notes and section.notes:
        body = f"{body}\n[Notes] {section.notes}" if body else f"[Notes] {section.notes}"
    section_text = format_section_chunk_text(
        heading=section.heading,
        body=body,
        page_number=section.page_number,
        doc_format=doc_format,
    )
    paragraphs = _split_into_paragraphs(section_text)
    if not paragraphs and section_text.strip():
        paragraphs = [section_text.strip()]
    return paragraphs


def _chunk_prose(
    parsed: ParsedDocument,
    max_tokens: int,
    overlap_tokens: int,
    doc_format: str,
) -> tuple[list[Chunk], dict[int, Section]]:
    """Paragraph-boundary chunker (legacy path; byte-identical at overlap 0)."""
    chunks: list[Chunk] = []
    origin: dict[int, Section] = {}
    idx = 0

    for section in parsed.sections:
        paragraphs = _section_paragraphs(section, doc_format, fold_notes=True)
        if not paragraphs:
            continue
        for piece in _pack_paragraphs(paragraphs, max_tokens, overlap_tokens):
            chunks.append(
                Chunk(
                    text=piece,
                    chunk_index=idx,
                    source_page=section.page_number,
                    section_heading=section.heading,
                    token_count=_token_count(piece),
                    doc_format=doc_format,
                )
            )
            origin[idx] = section
            idx += 1

    return chunks, origin


def _slide_title_base(title: str) -> str:
    """Normalize a slide title for continuation grouping (CH.3)."""
    base = (title or "").strip()
    base = _CONT_SUFFIX_RE.sub("", base)
    base = _TRAILING_NUM_RE.sub("", base)
    return base.strip().lower()


def _assign_topic_units(sections: list[Section]) -> list[int]:
    """Map each slide index to a topic-unit id (CH.3).

    Consecutive slides whose normalized titles match (e.g. a repeated
    section header or a ``(cont'd)`` title) share a unit; a slide with no
    title always starts a new unit (we never merge untitled slides).
    """
    units: list[int] = []
    current = 0
    prev_base: str | None = None
    for index, section in enumerate(sections):
        base = _slide_title_base(section.heading)
        if index == 0:
            units.append(0)
        elif base and prev_base and base == prev_base:
            units.append(current)
        else:
            current += 1
            units.append(current)
        prev_base = base or None
    return units


def _chunk_deck(
    parsed: ParsedDocument,
    max_tokens: int,
    overlap_tokens: int,
    doc_format: str,
) -> tuple[list[Chunk], dict[int, Section]]:
    """Slide-boundary-preserving deck chunker (CH.2/CH.3).

    Guarantees: two slides are never merged into one chunk; a slide is
    split only when it exceeds ``max_tokens`` (recorded via ``slide_part`` /
    ``slide_parts``); speaker notes become a distinct ``chunk_role="notes"``
    chunk linked to the slide; continuation slides share a ``topic_unit``.
    """
    chunks: list[Chunk] = []
    origin: dict[int, Section] = {}
    idx = 0
    topic_units = _assign_topic_units(parsed.sections)

    def emit(section: Section, pieces: list[str], role: str, unit: int) -> None:
        nonlocal idx
        parts = len(pieces)
        for part_index, piece in enumerate(pieces):
            chunks.append(
                Chunk(
                    text=piece,
                    chunk_index=idx,
                    source_page=section.page_number,
                    section_heading=section.heading,
                    token_count=_token_count(piece),
                    doc_format=doc_format,
                    chunk_role=role,
                    slide_part=part_index,
                    slide_parts=parts,
                    topic_unit=unit,
                )
            )
            origin[idx] = section
            idx += 1

    for section_index, section in enumerate(parsed.sections):
        unit = topic_units[section_index]

        # Slide body: never merged with another slide's content. Skip a
        # body chunk that would consist solely of the slide-context prefix
        # (a notes-only slide has no heading/body of its own -- its notes
        # are emitted below as a distinct chunk).
        has_body_content = bool(section.text.strip() or section.heading.strip())
        body_paragraphs = _section_paragraphs(section, doc_format, fold_notes=False)
        if body_paragraphs and has_body_content:
            body_pieces = _pack_paragraphs(body_paragraphs, max_tokens, overlap_tokens)
            emit(section, body_pieces, "body", unit)

        # Speaker notes: distinct chunk(s) linked to the same slide + unit.
        if section.notes:
            notes_text = f"[Notes] {section.notes}"
            notes_paragraphs = _split_into_paragraphs(notes_text) or [notes_text]
            notes_pieces = _pack_paragraphs(notes_paragraphs, max_tokens, overlap_tokens)
            emit(section, notes_pieces, "notes", unit)

    return chunks, origin


_VISUAL_BODY_PREFIXES = (
    "[Visual omitted:",
    "[Visual (alt text):",
    "[Visual (caption):",
    "[Scanned",
)
_VISUAL_MARKER_PREFIXES = (*_VISUAL_BODY_PREFIXES, "[Slide ", "[Page ")


def _classify_chunk_kind(text: str) -> str:
    """Return ``"visual"`` / ``"text"`` / ``"mixed"`` for a chunk's text."""
    lines = [ln for ln in (s.strip() for s in text.splitlines()) if ln]
    if not lines:
        return "text"
    visual_lines = sum(1 for ln in lines if ln.startswith(_VISUAL_BODY_PREFIXES))
    body_lines = sum(1 for ln in lines if not ln.startswith(_VISUAL_MARKER_PREFIXES))
    if visual_lines and body_lines == 0:
        return "visual"
    if visual_lines:
        return "mixed"
    return "text"


def _backfill_visual_context(
    parsed: ParsedDocument,
    chunks: list[Chunk],
    origin: dict[int, Section],
) -> None:
    """Populate ``chunk_kind`` and ``visual_assets`` after the main loop."""
    assets = parsed.visual_diagnostics.assets
    for chunk in chunks:
        chunk.chunk_kind = _classify_chunk_kind(chunk.text)
        section = origin.get(chunk.chunk_index)
        if section is None or not section.visual_asset_indexes:
            continue
        chunk.visual_assets = [
            assets[i].to_dict() for i in section.visual_asset_indexes if 0 <= i < len(assets)
        ]
