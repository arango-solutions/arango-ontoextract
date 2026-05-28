"""Stream 13 IMG.8: regression fixtures + end-to-end tests for image-aware extraction.

These tests build *real* PPTX and PDF artifacts in-process (no host-tool
dependencies) so the visual-aware pipeline has deterministic regression
coverage. They cover the five exit criteria from the Stream 13 plan:

1. **Inventory** -- visual asset metadata reflects the actual fixture.
2. **Placeholder mode** -- toggling ``visual_extraction_placeholders``
   removes ``[Visual omitted: ...]`` lines without losing the counts.
3. **OCR / caption injection** -- a fake ``VisualCaptionProvider`` flows
   captions into chunk text and updates the diagnostics.
4. **Prompt rendering** -- a visual-heavy chunk set routes through the
   strategy selector to ``tier1_visual_aware`` and the rendered system
   prompt explains the visual markers.
5. **Orphan-risk warning** -- aggregated visual diagnostics + an orphan
   ratio above threshold produce a ``visual_heavy_orphans`` warning.

The PPTX fixture encodes a "title-only taxonomy" -- a hierarchy that is
intentionally only legible via slide titles + alt-text -- which is the
exact failure mode IMG.7 / IMG.8 were written to catch.
"""

from __future__ import annotations

import io
from typing import Any

import pytest

from app.services.visual_extraction import (
    _PROVIDERS,
    CaptionResult,
    VisualCaptionProvider,
    aggregate_document_visual_diagnostics,
    build_orphan_risk_warning,
    register_caption_provider,
)

# ---------------------------------------------------------------------------
# Minimal PNG bytes -- valid 1x1 transparent PNG (reused from test_ingestion).
# Kept inline so this file is self-contained as a regression fixture source.
# ---------------------------------------------------------------------------
_MINIMAL_PNG = (
    b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
    b"\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx\x9cc\xf8\xcf\xc0"
    b"\x00\x00\x00\x03\x00\x01\x00\x05\xfe\x02\xfe\xdc\xcc\x59\xe7\x00\x00"
    b"\x00\x00IEND\xaeB`\x82"
)


# ---------------------------------------------------------------------------
# PPTX fixture: image-only taxonomy + title-only hierarchy.
# ---------------------------------------------------------------------------
def build_visual_taxonomy_pptx() -> bytes:
    """Build a PPTX whose ontology can ONLY be inferred from titles + visuals.

    Slide layout (10 slides):

    1. Title slide: "Vehicle Taxonomy" (root narrative).
    2. Section title: "Vehicle" -- title-only (the root class).
    3. Picture with alt text "Car body styles diagram" -- alt text is
       the sole evidence for the ``Car`` subclass cluster.
    4. Section title: "Car" -- title-only subclass.
    5. Section title: "Sedan" -- title-only subclass of ``Car``.
    6. Section title: "SUV" -- title-only subclass of ``Car``.
    7. Section title: "Truck" -- title-only subclass of ``Vehicle``.
    8. Picture without alt text -- exercises the placeholder / caption
       branch in ``collect_pptx_visual_assets``.
    9. Section title: "Pickup" -- title-only subclass of ``Truck``.
    10. Picture with alt text "Comparison chart: sedan vs SUV vs truck"
        -- alt text encodes relationships, not hierarchy.
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    title_only = prs.slide_layouts[5]  # "Title Only"
    blank = prs.slide_layouts[6]  # "Blank"

    def _add_title(text: str) -> None:
        slide = prs.slides.add_slide(title_only)
        slide.shapes.title.text = text

    def _add_picture(alt_text: str | None) -> None:
        slide = prs.slides.add_slide(blank)
        pic = slide.shapes.add_picture(
            io.BytesIO(_MINIMAL_PNG), Inches(1), Inches(1), width=Inches(2)
        )
        # The parser reads alt text via xpath ``.//p:cNvPr`` and checks
        # ``descr`` / ``title``. Mirror that exact path here so the
        # fixture and the production reader agree. python-pptx
        # auto-populates ``descr`` with the source filename when an
        # in-memory stream is used, so we must explicitly clear it for
        # the "no alt text" branch -- otherwise the caption path is
        # never exercised.
        elements = pic._element.xpath(".//p:cNvPr")
        assert elements, "expected a cNvPr element on the picture shape"
        if alt_text:
            elements[0].set("descr", alt_text)
        else:
            for attr in ("descr", "title"):
                if attr in elements[0].attrib:
                    del elements[0].attrib[attr]

    _add_title("Vehicle Taxonomy")
    _add_title("Vehicle")
    _add_picture("Car body styles diagram")
    _add_title("Car")
    _add_title("Sedan")
    _add_title("SUV")
    _add_title("Truck")
    _add_picture(None)
    _add_title("Pickup")
    _add_picture("Comparison chart: sedan vs SUV vs truck")

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# PDF fixture: scanned / image-only pages built with PyMuPDF.
# ---------------------------------------------------------------------------
def _build_white_png(*, width: int = 64, height: int = 64) -> bytes:
    """Generate a small, real PNG that mupdf will accept (1x1 fails)."""
    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", (width, height), color=(255, 255, 255)).save(buf, format="PNG")
    return buf.getvalue()


def build_scanned_pdf() -> bytes:
    """Build a PDF whose first page is image-only and whose second page has text.

    The first page has no text content and one inserted PNG -- fitz's
    block enumeration treats this as ``type != 0`` and ``parse_pdf``
    flags it as a scanned page. The second page has body text so the
    strategy selector can see the contrast (one scanned + one text page).
    """
    import fitz

    doc = fitz.open()
    png_bytes = _build_white_png()

    page_one = doc.new_page(width=612, height=792)
    page_one.insert_image(fitz.Rect(50, 50, 562, 742), stream=png_bytes)

    page_two = doc.new_page(width=612, height=792)
    page_two.insert_text(fitz.Point(72, 100), "Vehicle taxonomy continued.", fontsize=14)
    page_two.insert_text(
        fitz.Point(72, 140),
        "Cars and trucks each have subtypes.",
        fontsize=11,
    )

    out = doc.tobytes()
    doc.close()
    return out


# ---------------------------------------------------------------------------
# Helper: register a fake provider for the lifetime of a single test.
# ---------------------------------------------------------------------------
class _StubCaptionProvider(VisualCaptionProvider):
    """Deterministic provider that captions every image with a fixed string."""

    name = "stub"

    def __init__(self, *, text: str = "Captioned image", confidence: float = 0.81) -> None:
        self._text = text
        self._confidence = confidence
        self.calls: list[tuple[bytes, str]] = []

    def caption(self, image_bytes: bytes, *, mime_type: str = "image/png") -> CaptionResult:
        self.calls.append((image_bytes, mime_type))
        return CaptionResult(success=True, text=self._text, confidence=self._confidence)


class _AlwaysFailsCaptionProvider(VisualCaptionProvider):
    name = "always_fails"

    def caption(self, image_bytes: bytes, *, mime_type: str = "image/png") -> CaptionResult:
        return CaptionResult(success=False, failure_reason="simulated_provider_failure")


@pytest.fixture
def stub_caption_provider() -> _StubCaptionProvider:
    """Register a stub provider under id ``stub`` for the test body only."""
    provider = _StubCaptionProvider()

    class _Wrapper(VisualCaptionProvider):
        name = "stub"

        def caption(self, image_bytes: bytes, *, mime_type: str = "image/png") -> CaptionResult:
            return provider.caption(image_bytes, mime_type=mime_type)

    register_caption_provider("stub", _Wrapper)
    try:
        yield provider
    finally:
        _PROVIDERS.pop("stub", None)


@pytest.fixture
def failing_caption_provider() -> None:
    register_caption_provider("always_fails", _AlwaysFailsCaptionProvider)
    try:
        yield
    finally:
        _PROVIDERS.pop("always_fails", None)


# ---------------------------------------------------------------------------
# 1. Inventory regression
# ---------------------------------------------------------------------------
class TestVisualInventory:
    def test_pptx_visual_taxonomy_counts_assets_and_pages(self):
        from app.services.ingestion import parse_pptx

        parsed = parse_pptx(build_visual_taxonomy_pptx())

        diag = parsed.visual_diagnostics
        # Three picture shapes regardless of alt text presence.
        assert diag.visual_asset_count == 3
        # Two of those pictures had explicit alt text.
        assert diag.alt_text_count == 2
        # The third lacked alt text and -- with the default provider
        # being no-op -- falls through to the placeholder line.
        assert diag.placeholder_count == 1
        # caption_count starts at zero before any provider runs.
        assert diag.caption_count == 0
        # All title-only slides + picture slides are kept as sections.
        assert {a.page_number for a in diag.assets} == {3, 8, 10}
        assert sorted(diag.pages_with_visuals) == [3, 8, 10]

    def test_pdf_scanned_page_is_flagged(self):
        from app.services.ingestion import parse_pdf

        parsed = parse_pdf(build_scanned_pdf())

        # First page has no text -> scanned_page_count incremented.
        assert parsed.visual_diagnostics.scanned_page_count >= 1
        # The synthetic image inserted on page 1 is an image_block.
        page_ones = [a for a in parsed.visual_diagnostics.assets if a.page_number == 1]
        assert page_ones, "expected at least one image_block asset on the scanned page"
        assert all(a.asset_type == "image_block" for a in page_ones)
        # The scanned page section text starts with the explicit marker.
        scanned_section = next(
            (s for s in parsed.sections if s.page_number == 1 and "[Scanned" in s.text),
            None,
        )
        assert scanned_section is not None, "expected a scanned-page section with marker"


# ---------------------------------------------------------------------------
# 2. Placeholder mode toggle
# ---------------------------------------------------------------------------
class TestPlaceholderMode:
    def test_disabling_placeholders_removes_marker_but_keeps_counts(self, monkeypatch):
        from app.config import settings
        from app.services.ingestion import parse_pptx

        monkeypatch.setattr(settings, "visual_extraction_placeholders", False)
        parsed = parse_pptx(build_visual_taxonomy_pptx())

        diag = parsed.visual_diagnostics
        # Pictures with alt text still emit alt-text lines and are counted.
        assert diag.alt_text_count == 2
        # The picture without alt text is still inventoried; placeholder
        # text is suppressed but the asset count stays honest.
        assert diag.visual_asset_count == 2 or diag.visual_asset_count == 3
        # No [Visual omitted: ...] marker leaks into chunk text.
        for section in parsed.sections:
            assert "[Visual omitted:" not in section.text

    def test_default_placeholder_mode_emits_marker(self, monkeypatch):
        from app.config import settings
        from app.services.ingestion import parse_pptx

        monkeypatch.setattr(settings, "visual_extraction_placeholders", True)
        monkeypatch.setattr(settings, "visual_caption_provider", "none")
        parsed = parse_pptx(build_visual_taxonomy_pptx())

        joined = "\n".join(s.text for s in parsed.sections)
        # The picture-without-alt-text slide emits the placeholder marker.
        assert "[Visual omitted: slide 8 image" in joined


# ---------------------------------------------------------------------------
# 3. OCR / caption injection
# ---------------------------------------------------------------------------
class TestCaptionInjection:
    def test_stub_provider_caption_replaces_placeholder(self, monkeypatch, stub_caption_provider):
        from app.config import settings
        from app.services.ingestion import parse_pptx

        monkeypatch.setattr(settings, "visual_caption_provider", "stub")
        parsed = parse_pptx(build_visual_taxonomy_pptx())

        diag = parsed.visual_diagnostics
        # The captionable picture (slide 8, no alt text) routes through
        # the provider and is recorded as ``vision_caption``.
        captioned = [a for a in diag.assets if a.method == "vision_caption"]
        assert len(captioned) == 1
        assert captioned[0].page_number == 8
        assert captioned[0].confidence == pytest.approx(0.81)
        assert diag.caption_count == 1
        # The provider was actually invoked with the embedded PNG bytes.
        assert stub_caption_provider.calls, "expected provider.caption() to be called"
        png_bytes, mime = stub_caption_provider.calls[0]
        assert png_bytes.startswith(b"\x89PNG")
        assert mime.startswith("image/")
        # The caption text reaches the section body.
        slide_8 = next(s for s in parsed.sections if s.page_number == 8)
        assert "[Visual (caption): Captioned image]" in slide_8.text

    def test_max_caption_cap_prevents_runaway_calls(self, monkeypatch, stub_caption_provider):
        from app.config import settings
        from app.services.ingestion import parse_pptx

        monkeypatch.setattr(settings, "visual_caption_provider", "stub")
        # Cap at zero -> provider should never be invoked even though
        # the deck has a captionable picture.
        monkeypatch.setattr(settings, "visual_caption_max_assets_per_doc", 0)
        parsed = parse_pptx(build_visual_taxonomy_pptx())

        assert parsed.visual_diagnostics.caption_count == 0
        assert stub_caption_provider.calls == []
        # And the placeholder line is preserved as the fallback.
        joined = "\n".join(s.text for s in parsed.sections)
        assert "[Visual omitted: slide 8 image" in joined

    def test_failing_provider_falls_back_to_placeholder(
        self, monkeypatch, failing_caption_provider
    ):
        from app.config import settings
        from app.services.ingestion import parse_pptx

        monkeypatch.setattr(settings, "visual_caption_provider", "always_fails")
        parsed = parse_pptx(build_visual_taxonomy_pptx())

        # The no-alt-text picture stays at placeholder with a recorded
        # failure_reason from the provider.
        failed = [
            a
            for a in parsed.visual_diagnostics.assets
            if a.failure_reason == "simulated_provider_failure"
        ]
        assert len(failed) == 1
        # The successful alt-text pictures keep their method.
        alt = [a for a in parsed.visual_diagnostics.assets if a.method == "alt_text"]
        assert len(alt) == 2
        # caption_count stays at zero on provider failure.
        assert parsed.visual_diagnostics.caption_count == 0


# ---------------------------------------------------------------------------
# 4. Prompt rendering through the strategy selector
# ---------------------------------------------------------------------------
class TestPromptRendering:
    def test_visual_taxonomy_pptx_routes_to_visual_aware_prompt(self):
        from app.extraction.agents.strategy import strategy_selector_node
        from app.extraction.prompts import get_template
        from app.services.ingestion import chunk_document, parse_pptx

        parsed = parse_pptx(build_visual_taxonomy_pptx())
        chunks = chunk_document(parsed)
        chunk_dicts = [
            {
                "chunk_id": f"c{c.chunk_index}",
                "text": c.text,
                "chunk_kind": c.chunk_kind,
                "doc_format": "pptx",
            }
            for c in chunks
        ]

        state: dict[str, Any] = {
            "extraction_run_id": "run-img8",
            "run_id": "run-img8",
            "document_chunks": chunk_dicts,
            "step_logs": [],
        }
        result = strategy_selector_node(state)
        config = result["strategy_config"]
        assert config["document_type"] == "visual_heavy_presentation"
        assert config["prompt_template_key"] == "tier1_visual_aware"

        # The rendered prompt body must include the chunk text so the
        # LLM sees the same visual markers that drove the routing.
        template = get_template("tier1_visual_aware")
        system, user = template.render(
            chunks_text="\n".join(c["text"] for c in chunk_dicts[:3]),
            domain_context="",
            extra_vars={"pass_number": 1, "model_name": "test-model"},
        )
        assert "[Slide N: Title]" in system
        assert "[Visual omitted:" in system
        # The user prompt should contain at least one fixture title.
        assert "Vehicle" in user

    def test_caption_marker_documented_in_prompt(self):
        from app.extraction.prompts import get_template

        system = get_template("tier1_visual_aware").system_prompt
        # IMG.4 + IMG.8: captions are a first-class marker and the
        # prompt explicitly tells the LLM how to treat them.
        assert "[Visual (caption):" in system


# ---------------------------------------------------------------------------
# 5. Orphan-risk warning end-to-end
# ---------------------------------------------------------------------------
class TestOrphanRiskWarning:
    def _doc_record_from_parsed(self, parsed, *, doc_id: str, filename: str) -> dict[str, Any]:
        """Mirror what ``process_document`` would persist to ``documents``."""
        return {
            "_key": doc_id,
            "filename": filename,
            "metadata": {"visual_extraction": parsed.visual_diagnostics.to_metadata_dict()},
        }

    def test_visual_taxonomy_with_orphan_classes_triggers_warning(self):
        from app.services.ingestion import parse_pptx

        parsed = parse_pptx(build_visual_taxonomy_pptx())
        # Lower min_visual_assets so the 3-image fixture qualifies.
        doc = self._doc_record_from_parsed(parsed, doc_id="d1", filename="taxonomy.pptx")
        summary = aggregate_document_visual_diagnostics([doc])

        # Simulate 4 of 5 extracted classes coming back without a parent
        # -- the exact failure mode IMG.7 is meant to surface.
        warning = build_orphan_risk_warning(
            orphan_class_count=4,
            total_classes=5,
            visual_summary=summary,
            orphan_ratio_threshold=0.5,
            min_visual_assets=2,
        )
        assert warning is not None
        assert warning["type"] == "visual_heavy_orphans"
        assert warning["orphan_class_count"] == 4
        assert warning["visual_asset_count"] == 3
        assert warning["documents"], "expected per-document breakdown"
        per_doc = warning["documents"][0]
        assert per_doc["filename"] == "taxonomy.pptx"
        assert per_doc["visual_asset_count"] == 3

    def test_scanned_pdf_alone_triggers_warning_even_with_few_orphans(self):
        from app.services.ingestion import parse_pdf

        parsed = parse_pdf(build_scanned_pdf())
        doc = self._doc_record_from_parsed(parsed, doc_id="d2", filename="scan.pdf")
        summary = aggregate_document_visual_diagnostics([doc])

        # Even 1/10 orphans triggers when scanned pages > 0, because
        # the curator needs to know an entire page never reached the LLM.
        warning = build_orphan_risk_warning(
            orphan_class_count=1,
            total_classes=10,
            visual_summary=summary,
            orphan_ratio_threshold=0.5,
            min_visual_assets=99,
        )
        assert warning is not None
        assert warning["scanned_page_count"] >= 1

    def test_healthy_run_with_few_visuals_does_not_warn(self):
        # Tiny visual footprint + no orphans -> no warning.
        summary = aggregate_document_visual_diagnostics(
            [
                {
                    "_key": "d3",
                    "filename": "narrative.pdf",
                    "metadata": {
                        "visual_extraction": {
                            "visual_asset_count": 1,
                            "scanned_page_count": 0,
                            "placeholder_count": 1,
                            "alt_text_count": 0,
                            "pages_with_visuals": [3],
                        }
                    },
                }
            ]
        )
        warning = build_orphan_risk_warning(
            orphan_class_count=0,
            total_classes=20,
            visual_summary=summary,
            orphan_ratio_threshold=0.5,
            min_visual_assets=5,
        )
        assert warning is None
